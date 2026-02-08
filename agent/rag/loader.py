"""文档加载器 — 支持多种格式

文本类（直接提取文字）:
  .txt, .md, .pdf, .docx, .xlsx, .pptx

多媒体类（调 DashScope API 转文字）:
  图片: .png, .jpg, .jpeg, .bmp, .webp, .gif
  音频: .mp3, .wav, .flac, .m4a, .ogg, .aac
  视频: .mp4, .avi, .mov, .mkv, .webm
"""

import base64
import json
import mimetypes
import signal
import threading

from typing import List, Optional

from langchain_core.documents import Document

from agent.config import DASHSCOPE_API_KEY, VISION_MODEL, ASR_MODEL
from pathlib import Path

def iter_supported_files(path: str):
    p = Path(path)
    if p.is_file():
        yield p
        return
    if not p.is_dir():
        raise FileNotFoundError(f"路径不存在: {path}")

    for f in p.rglob("*"):
        if f.is_file() and f.suffix.lower() in {".txt", ".md", ".pdf"}:
            yield f


# ──────────────────────── 超时工具 ────────────────────────

class _TimeoutError(Exception):
    """API 调用超时。"""


def _run_with_timeout(func, args=(), timeout=60):
    """在子线程中运行函数，超时则放弃（Windows 兼容）。

    用线程而非 signal.alarm，因为 Windows 不支持 SIGALRM。
    """
    result = [None]
    error = [None]

    def worker():
        try:
            result[0] = func(*args)
        except Exception as e:
            error[0] = e

    t = threading.Thread(target=worker, daemon=True)
    t.start()
    t.join(timeout=timeout)

    if t.is_alive():
        # 线程还在跑 = 超时，daemon 线程会在主进程退出时被清理
        raise _TimeoutError(f"超时 ({timeout}s)")
    if error[0]:
        raise error[0]
    return result[0]

# ──────────────────────── 文本类加载器 ────────────────────────


def _load_txt(path: str) -> List[Document]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": path})]


def _load_md(path: str) -> List[Document]:
    return _load_txt(path)


def _load_pdf(path: str) -> List[Document]:
    try:
        from langchain_community.document_loaders import PyPDFLoader
        return PyPDFLoader(path).load()
    except ImportError:
        raise ImportError("需要安装 pypdf: pip install pypdf")


def _load_docx(path: str) -> List[Document]:
    from docx import Document as DocxDocument

    doc = DocxDocument(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    return [Document(page_content=text, metadata={"source": path})]


def _load_xlsx(path: str) -> List[Document]:
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True)
    docs = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
            docs.append(Document(
                page_content=text,
                metadata={"source": path, "sheet": sheet_name},
            ))
    return docs


def _load_pptx(path: str) -> List[Document]:
    from pptx import Presentation

    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            text = f"[Slide {i}]\n" + "\n".join(texts)
            docs.append(Document(
                page_content=text,
                metadata={"source": path, "slide": i},
            ))
    return docs


# ──────────────────────── 多媒体加载器 ────────────────────────

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".gif"}
AUDIO_EXTS = {".mp3", ".wav", ".flac", ".m4a", ".ogg", ".aac"}
VIDEO_EXTS = {".mp4", ".avi", ".mov", ".mkv", ".webm"}


def _describe_image(path: str) -> List[Document]:
    """用 DashScope 视觉模型描述图片内容。"""
    import httpx

    with open(path, "rb") as f:
        img_bytes = f.read()
    b64 = base64.b64encode(img_bytes).decode()

    mime = mimetypes.guess_type(path)[0] or "image/png"

    resp = httpx.post(
        "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {DASHSCOPE_API_KEY}",
            "Content-Type": "application/json",
        },
        json={
            "model": VISION_MODEL,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{mime};base64,{b64}"},
                        },
                        {
                            "type": "text",
                            "text": "请详细描述这张图片的内容。如果图中有文字，请完整提取出来。",
                        },
                    ],
                }
            ],
        },
        timeout=60.0,
    )
    resp.raise_for_status()
    text = resp.json()["choices"][0]["message"]["content"]
    return [Document(page_content=text, metadata={"source": path, "type": "image"})]


def _transcribe_audio(path: str) -> List[Document]:
    """用 DashScope 语音识别模型转录音频。"""
    import dashscope
    from dashscope.audio.asr import Transcription

    dashscope.api_key = DASHSCOPE_API_KEY

    # 上传文件获取 URL
    file_urls = [f"file://{Path(path).resolve()}"]

    result = Transcription.call(
        model=ASR_MODEL,
        file_urls=file_urls,
    )

    if result.status_code != 200:
        return [Document(
            page_content=f"[音频转录失败: {result.message}]",
            metadata={"source": path, "type": "audio"},
        )]

    # 提取转录文本
    texts = []
    if hasattr(result, "output") and result.output:
        results = result.output.get("results", [])
        for r in results:
            transcript = r.get("transcription_url") or r.get("text", "")
            if isinstance(transcript, str) and transcript:
                texts.append(transcript)

    text = "\n".join(texts) if texts else "[未识别到音频内容]"
    return [Document(page_content=text, metadata={"source": path, "type": "audio"})]


def _load_video(path: str) -> List[Document]:
    """处理视频：用 DashScope 视觉模型理解视频内容。"""
    import httpx

    # DashScope qwen-vl-max 支持直接传视频文件 (base64)
    # 但视频文件可能很大，优先用视频 URL 或截取帧
    # 这里用 DashScope 的视频理解能力 (qwen-vl-max 支持 video_url)
    file_path = Path(path).resolve()
    file_size_mb = file_path.stat().st_size / (1024 * 1024)

    if file_size_mb > 50:
        return [Document(
            page_content=f"[视频文件过大 ({file_size_mb:.1f}MB)，跳过处理]",
            metadata={"source": path, "type": "video"},
        )]

    with open(path, "rb") as f:
        video_bytes = f.read()
    b64 = base64.b64encode(video_bytes).decode()

    mime = mimetypes.guess_type(path)[0] or "video/mp4"

    resp = httpx.post(
        "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {DASHSCOPE_API_KEY}",
            "Content-Type": "application/json",
        },
        json={
            "model": VISION_MODEL,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "video_url",
                            "video_url": {"url": f"data:{mime};base64,{b64}"},
                        },
                        {
                            "type": "text",
                            "text": "请详细描述这个视频的内容，包括画面中出现的文字、人物动作、场景变化等关键信息。",
                        },
                    ],
                }
            ],
        },
        timeout=120.0,
    )
    resp.raise_for_status()
    text = resp.json()["choices"][0]["message"]["content"]
    return [Document(page_content=text, metadata={"source": path, "type": "video"})]


# ──────────────────────── 注册表 ────────────────────────

# 文本类
LOADERS = {
    ".txt": _load_txt,
    ".md": _load_md,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".xlsx": _load_xlsx,
    ".pptx": _load_pptx,
}

# 多媒体类
for ext in IMAGE_EXTS:
    LOADERS[ext] = _describe_image
for ext in AUDIO_EXTS:
    LOADERS[ext] = _transcribe_audio
for ext in VIDEO_EXTS:
    LOADERS[ext] = _load_video


# ──────────────────────── 主入口 ────────────────────────

ALL_SUPPORTED_EXTS = set(LOADERS.keys())


def load_documents(path: str) -> List[Document]:
    """加载单个文件或整个目录下的文档。

    支持的格式: txt, md, pdf, docx, xlsx, pptx,
               png/jpg/bmp/webp/gif (图片→描述),
               mp3/wav/flac/m4a (音频→转录),
               mp4/avi/mov/mkv/webm (视频→描述)
    """
    p = Path(path)
    docs = []

    if p.is_file():
        files = [p]
    elif p.is_dir():
        files = sorted(f for f in p.rglob("*") if f.suffix.lower() in LOADERS)
    else:
        raise FileNotFoundError(f"路径不存在: {path}")

    total = len(files)
    for idx, f in enumerate(files, 1):
        ext = f.suffix.lower()
        loader = LOADERS.get(ext)
        if not loader:
            continue

        # 判断类型标签
        if ext in IMAGE_EXTS:
            tag = "图片→描述"
        elif ext in AUDIO_EXTS:
            tag = "音频→转录"
        elif ext in VIDEO_EXTS:
            tag = "视频→描述"
        else:
            tag = "文本"

        try:
            loaded = loader(str(f))
            print(f"  [{idx}/{total}] {f.name} ({tag}, {len(loaded)} 段)")
            docs.extend(loaded)
        except Exception as e:
            print(f"  [{idx}/{total}] {f.name} (失败: {e})")

    if not docs:
        print(f"  警告: 未在 {path} 中找到支持的文档")

    return docs
